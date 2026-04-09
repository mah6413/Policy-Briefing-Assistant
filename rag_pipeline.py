import math
import os
import re
from pathlib import Path
from typing import Any, Dict, List

import openpyxl
from docx import Document
from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation
from pypdf import PdfReader

load_dotenv()

OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
client = OpenAI(api_key=OPENAI_API_KEY) if OPENAI_API_KEY else None


def extract_text_from_pdf(file_path: str) -> str:
    reader = PdfReader(file_path)
    pages = []
    for page in reader.pages:
        pages.append(page.extract_text() or "")
    return "\n".join(pages).strip()


def extract_text_from_docx(file_path: str) -> str:
    doc = Document(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs).strip()


def extract_text_from_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text_runs.append(shape.text)
    return "\n".join(text_runs).strip()


def extract_text_from_xlsx(file_path: str) -> str:
    workbook = openpyxl.load_workbook(file_path, data_only=True)
    rows_text = []

    for sheet in workbook.worksheets:
        rows_text.append(f"Sheet: {sheet.title}")
        for row in sheet.iter_rows(values_only=True):
            values = [str(cell) for cell in row if cell is not None]
            if values:
                rows_text.append(" | ".join(values))

    return "\n".join(rows_text).strip()


def extract_text_by_file_type(file_path: str) -> str:
    ext = Path(file_path).suffix.lower()

    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    if ext == ".docx":
        return extract_text_from_docx(file_path)
    if ext == ".pptx":
        return extract_text_from_pptx(file_path)
    if ext == ".xlsx":
        return extract_text_from_xlsx(file_path)

    raise ValueError(f"Unsupported file type: {ext}")


def build_document_library(uploaded_file_paths: List[str]) -> List[Dict[str, Any]]:
    document_library = []

    for i, file_path in enumerate(uploaded_file_paths, start=1):
        file_name = Path(file_path).name
        ext = Path(file_path).suffix.lower().replace(".", "")
        extracted_text = extract_text_by_file_type(file_path)

        document_library.append(
            {
                "document_id": f"doc_{i}",
                "document_title": file_name,
                "document_type": ext,
                "text": extracted_text,
            }
        )

    return document_library


def chunk_text(text: str, chunk_size: int = 1200, overlap: int = 150) -> List[str]:
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return []

    chunks = []
    start = 0
    text_length = len(text)

    while start < text_length:
        end = min(start + chunk_size, text_length)
        chunks.append(text[start:end])

        if end == text_length:
            break

        start = max(0, end - overlap)

    return chunks


def build_chunk_store(document_library: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    chunk_store = []

    for doc in document_library:
        chunks = chunk_text(doc["text"])
        for idx, chunk in enumerate(chunks, start=1):
            chunk_store.append(
                {
                    "chunk_id": f"{doc['document_id']}_chunk_{idx}",
                    "document_id": doc["document_id"],
                    "document_title": doc["document_title"],
                    "document_type": doc["document_type"],
                    "text": chunk,
                    "embedding": None,
                }
            )

    return chunk_store


def get_embedding(text: str) -> List[float]:
    if client is None:
        raise ValueError("OPENAI_API_KEY is not set.")

    response = client.embeddings.create(
        model="text-embedding-3-small",
        input=text,
    )
    return response.data[0].embedding


def embed_chunk_store(chunk_store: List[Dict[str, Any]]) -> List[Dict[str, Any]]:
    for chunk in chunk_store:
        chunk["embedding"] = get_embedding(chunk["text"])
    return chunk_store


def cosine_similarity(vec_a: List[float], vec_b: List[float]) -> float:
    dot = sum(a * b for a, b in zip(vec_a, vec_b))
    norm_a = math.sqrt(sum(a * a for a in vec_a))
    norm_b = math.sqrt(sum(b * b for b in vec_b))

    if norm_a == 0 or norm_b == 0:
        return 0.0

    return dot / (norm_a * norm_b)


def retrieve_relevant_chunks(
    user_query: str, chunk_store: List[Dict[str, Any]], top_k: int = 6
) -> List[Dict[str, Any]]:
    query_embedding = get_embedding(user_query)

    scored_chunks = []
    for chunk in chunk_store:
        if chunk["embedding"] is None:
            continue
        score = cosine_similarity(query_embedding, chunk["embedding"])
        scored_chunk = {**chunk, "score": round(score, 4)}
        scored_chunks.append(scored_chunk)

    scored_chunks.sort(key=lambda x: x["score"], reverse=True)
    return scored_chunks[:top_k]


def evaluate_retrieval_quality(sources: List[Dict[str, Any]]) -> str:
    if not sources:
        return "no_relevant_context"

    top_score = sources[0].get("score", 0)

    if top_score >= 0.75:
        return "strong"
    if top_score >= 0.55:
        return "moderate"
    return "no_relevant_context"


def call_chat_completion(system_prompt: str, user_prompt: str) -> str:
    if client is None:
        raise ValueError("OPENAI_API_KEY is not set.")

    response = client.chat.completions.create(
        model="gpt-4.1-mini",
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        temperature=0.2,
    )
    return response.choices[0].message.content.strip()


def answer_question(user_query: str, chunk_store: List[Dict[str, Any]]) -> Dict[str, Any]:
    sources = retrieve_relevant_chunks(user_query, chunk_store, top_k=6)
    quality = evaluate_retrieval_quality(sources)

    context = "\n\n".join([f"[{s['document_title']}]\n{s['text']}" for s in sources])

    system_prompt = (
        "You are an AI Policy Briefing Assistant. "
        "Answer only from the provided document context. "
        "If the answer is not supported by the context, say so clearly."
    )

    user_prompt = f"""
Question:
{user_query}

Context:
{context}
"""

    answer = call_chat_completion(system_prompt, user_prompt)

    return {
        "answer": answer,
        "sources": sources,
        "confidence": "high" if quality == "strong" else "medium",
        "fallback_used": quality == "no_relevant_context",
    }


def generate_policy_memo(user_query: str, chunk_store: List[Dict[str, Any]]) -> str:
    sources = retrieve_relevant_chunks(user_query, chunk_store, top_k=6)
    context = "\n\n".join(f"[{s['document_title']}]\n{s['text']}" for s in sources)

    system_prompt = (
        "You are an AI Policy Briefing Assistant. "
        "Write a professional policy memo grounded only in the provided context."
    )

    user_prompt = f"""
Draft a policy memo on the following request:
{user_query}

Context:
{context}
"""

    return call_chat_completion(system_prompt, user_prompt)


def generate_policy_brief(user_query: str, chunk_store: List[Dict[str, Any]]) -> str:
    sources = retrieve_relevant_chunks(user_query, chunk_store, top_k=6)
    context = "\n\n".join(f"[{s['document_title']}]\n{s['text']}" for s in sources)

    system_prompt = (
        "You are an AI Policy Briefing Assistant. "
        "Write a concise policy brief grounded only in the provided context."
    )

    user_prompt = f"""
Draft a policy brief on the following request:
{user_query}

Context:
{context}
"""

    return call_chat_completion(system_prompt, user_prompt)


def generate_legislative_letter(user_query: str, chunk_store: List[Dict[str, Any]]) -> str:
    sources = retrieve_relevant_chunks(user_query, chunk_store, top_k=6)
    context = "\n\n".join(f"[{s['document_title']}]\n{s['text']}" for s in sources)

    system_prompt = (
        "You are an AI Policy Briefing Assistant. "
        "Draft a formal legislative letter grounded only in the provided context."
    )

    user_prompt = f"""
Draft a legislative letter for the following request:
{user_query}

Context:
{context}
"""

    return call_chat_completion(system_prompt, user_prompt)


def generate_executive_summary(user_query: str, chunk_store: List[Dict[str, Any]]) -> str:
    sources = retrieve_relevant_chunks(user_query, chunk_store, top_k=6)
    context = "\n\n".join(f"[{s['document_title']}]\n{s['text']}" for s in sources)

    system_prompt = (
        "You are an AI Policy Briefing Assistant. "
        "Write an executive summary grounded only in the provided context."
    )

    user_prompt = f"""
Draft an executive summary for the following request:
{user_query}

Context:
{context}
"""

    return call_chat_completion(system_prompt, user_prompt)


def clean_preview(text: str, max_chars: int = 700) -> str:
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return "No readable text could be extracted from the uploaded document."
    return text[:max_chars] + ("..." if len(text) > max_chars else "")


def bulletize_text(text: str, max_points: int = 3) -> List[str]:
    text = re.sub(r"\s+", " ", text).strip()
    if not text:
        return ["No clear findings could be extracted from the uploaded document."]

    sentences = re.split(r"(?<=[.!?])\s+", text)
    useful = [s.strip() for s in sentences if len(s.strip()) > 40]

    if not useful:
        useful = [text[:180] + ("..." if len(text) > 180 else "")]

    return useful[:max_points]


def get_demo_answer(
    mode: str, user_query: str, document_library: List[Dict[str, Any]]
) -> Dict[str, Any]:
    primary_doc = document_library[0] if document_library else None

    if primary_doc:
        doc_title = primary_doc["document_title"]
        doc_type = primary_doc["document_type"]
        doc_text = primary_doc["text"]
        preview = clean_preview(doc_text)
        findings = bulletize_text(doc_text, max_points=3)
    else:
        doc_title = "No document uploaded"
        doc_type = "unknown"
        preview = "No document content available."
        findings = ["No findings available."]

    findings_text = "\n".join([f"- {item}" for item in findings])

    templates = {
        "qa": f"""Demo Response

Question:
{user_query}

Sample grounded preview based on: {doc_title}

Document excerpt:
{preview}

This is a no-token demo response. In Live AI Mode, the assistant would retrieve relevant passages and generate a grounded answer using the uploaded document.
""",
        "policy_memo": f"""Policy Memo (Demo)

To: Leadership Team
From: AI Policy Briefing Assistant
Subject: {user_query}

Executive Summary:
This demonstration memo is based on a preview of the uploaded document, "{doc_title}". It shows how the assistant can turn source material into a structured memo without using OpenAI tokens in Demo Mode.

Background:
The uploaded file was identified as a {doc_type.upper()} document. A sample excerpt from the document is shown below to simulate grounded policy analysis.

Document Preview:
{preview}

Key Issues Identified:
{findings_text}

Recommendations:
- Review the uploaded document for the most decision-relevant sections
- Use Live AI Mode for a fully generated grounded memo
- Keep Demo Mode enabled for public-facing showcase use
""",
        "policy_brief": f"""Policy Brief (Demo)

Title: {user_query}

Source Document:
{doc_title}

Overview:
This is a demonstration policy brief generated from the uploaded file preview.

Document Preview:
{preview}

Illustrative Findings:
{findings_text}

Next Step:
Switch to Live AI Mode to generate a full document-grounded brief.
""",
        "legislative_letter": f"""Legislative Letter (Demo)

Dear [Legislator Name],

I am writing regarding {user_query}. This demonstration letter is based on a preview of the uploaded document "{doc_title}" and shows how the assistant can transform source material into advocacy-oriented writing.

Relevant Document Preview:
{preview}

In Live AI Mode, the assistant would generate a more complete and tailored draft grounded in the uploaded material.

Sincerely,
[Your Organization]
""",
        "executive_summary": f"""Executive Summary (Demo)

Topic:
{user_query}

Source:
{doc_title}

Summary:
This demonstration summary uses extracted text from the uploaded file to simulate a grounded output without calling the OpenAI API.

Document Preview:
{preview}

Illustrative Takeaways:
{findings_text}

Recommendation:
Use Live AI Mode for a fuller synthesis based on the uploaded document.
""",
    }

    return {
        "answer": templates.get(mode, templates["qa"]),
        "sources": [
            {
                "score": 1.0,
                "chunk_id": "demo_chunk_1",
                "document_id": primary_doc["document_id"] if primary_doc else "demo_doc",
                "document_title": doc_title,
                "document_type": doc_type,
                "text": preview,
            }
        ]
        if primary_doc
        else [],
        "confidence": "demo",
        "fallback_used": False,
    }


def run_policy_assistant(
    user_query: str,
    uploaded_file_paths: List[str],
    mode: str = "qa",
    demo_mode: bool = False,
) -> Dict[str, Any]:
    document_library = build_document_library(uploaded_file_paths)

    if demo_mode:
        return get_demo_answer(mode, user_query, document_library)

    if client is None:
        raise ValueError(
            "OPENAI_API_KEY is not set. Add it to your environment or enable Demo Mode."
        )

    chunk_store = build_chunk_store(document_library)
    chunk_store = embed_chunk_store(chunk_store)

    if mode == "qa":
        return answer_question(user_query, chunk_store)
    elif mode == "policy_memo":
        text = generate_policy_memo(user_query, chunk_store)
    elif mode == "policy_brief":
        text = generate_policy_brief(user_query, chunk_store)
    elif mode == "legislative_letter":
        text = generate_legislative_letter(user_query, chunk_store)
    elif mode == "executive_summary":
        text = generate_executive_summary(user_query, chunk_store)
    else:
        raise ValueError(f"Unsupported mode: {mode}")

    sources = retrieve_relevant_chunks(user_query, chunk_store, top_k=6)
    quality = evaluate_retrieval_quality(sources)

    return {
        "answer": text,
        "sources": sources,
        "confidence": "high" if quality == "strong" else "medium",
        "fallback_used": quality == "no_relevant_context",
    }